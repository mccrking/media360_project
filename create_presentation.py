#!/usr/bin/env python
"""Generate a comprehensive PowerPoint presentation for MediaPulse 360 project."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(13, 27, 42)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = ORANGE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "📰 Veille Automatisée - Architecture des Données"
    footer_frame.paragraphs[0].font.size = Pt(14)
    footer_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY

def add_content_slide(prs, title, content_points, use_bg=True):
    """Add content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    if use_bg:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0

# Slide 1: Title
add_title_slide(prs, "MediaPulse 360", "Dashboard de Veille Médias Automatisé")

# Slide 2: Objectives
add_content_slide(prs, "Objectifs du Projet", [
    "✓ Ingestion automatique d'articles depuis 4 sources (BBC, Reuters, Hespress, France24)",
    "✓ Stockage structuré en Data Warehouse PostgreSQL",
    "✓ Dashboard interactif en temps réel avec Streamlit",
    "✓ Quarantaine intelligente des articles suspects",
    "✓ Audit complet des changements d'état",
    "✓ Titres cliquables pour navigation directe"
])

# Slide 3: Architecture complète UPDATED
add_content_slide(prs, "Architecture du Système", [
    "🔄 **Bronze Layer** → RSS Scraper → MinIO raw articles (JSONL)",
    "🧹 **Silver Layer** → Nettoyage & normalisation → MinIO + PostgreSQL",
    "📊 **Gold Layer** → Agrégations analytiques → PostgreSQL tables",
    "🎯 **Orchestration** → Apache Airflow (DAGs, Scheduling, Monitoring)",
    "📺 **Visualisation** → Streamlit Dashboard",
    "🔐 **Data Quality** → Validation des 3 dimensions (Complétude, Cohérence, Validité)"
])

# Slide 4: Sources
add_content_slide(prs, "Sources de Données", [
    "• BBC News (https://feeds.bbci.co.uk/news/world/rss.xml)",
    "  → 22 articles actuellement en BD",
    "• Reuters World (reutersagency.com/feed)",
    "  → 2 articles actuellement en BD",
    "• Hespress (https://fr.hespress.com/feed)",
    "  → 15 articles actuellement en BD",
    "• France24 (https://www.france24.com/fr/rss) [NOUVEAU]",
    "  → 20 articles actuellement en BD"
])

# Slide 5: Key Features
add_content_slide(prs, "Fonctionnalités Clés", [
    "🎯 Filtrage par source en temps réel",
    "📅 Fenêtre d'analyse configurable (1-30 jours)",
    "🔐 Quarantaine avec raisons annotées",
    "✅ Approbation en masse des articles",
    "🔗 Canonical URLs pour stabilité des IDs",
    "📊 KPIs : Total articles, Source dominante",
    "📈 Graphiques : Distribution par source, Timeline"
])

# Slide 6: Data Model
add_content_slide(prs, "Modèle de Données", [
    "articles_detail:",
    "  • article_id (sha256 stable) | title | author | published_at",
    "  • category | content | source | url | canonical_url",
    "  • quarantine flag | scraped_at",
    "",
    "quarantine_audit (Audit complet):",
    "  • action (ingest, approve, reject) | performed_by | reason",
    "  • performed_at (TIMESTAMP)"
])

# Slide 7: Technology Stack UPDATED
add_content_slide(prs, "Stack Technologique", [
    "🐍 Python 3.11 → Scraping & Ingestion",
    "📊 Streamlit → Interface interactive",
    "🗄️ PostgreSQL 16 → Data Warehouse + Métadata",
    "☁️ MinIO → Data Lake distribuée (Bronze/Silver/Gold)",
    "🔄 Apache Airflow 2.8 → Orchestration & Scheduling",
    "🐳 Docker Compose → Déploiement et orchestration services"
])

# Slide 8: Current Status
add_content_slide(prs, "État Actuel", [
    "✅ Scraper → 50 articles par cycle",
    "✅ Ingestion → Upsert automatique + Quarantaine",
    "✅ Database → 59 articles uniques",
    "✅ Dashboard → Live à http://localhost:8501",
    "✅ Migrations → 3 versions (init, canonical/audit, backfill)",
    "✅ Auto-refresh → Toutes les 5 minutes",
    "🔄 Ingestion continue en Docker"
])

# Slide 9: Quarantine Heuristics
add_content_slide(prs, "Heuristiques de Quarantaine", [
    "Critères de flagging automatique:",
    "",
    "• Domaines sensibles (BBC, Hespress, Reuters)",
    "• Titres génériques ou court (<8 caractères)",
    "• Contenu fragmenté (<100 mots)",
    "• Patterns suspects en URL",
    "",
    "Raisons annotées → Audit trail complet"
])

# Slide 10: Dashboard Experience
add_content_slide(prs, "Expérience Utilisateur", [
    "🎨 Interface Dark Theme moderne",
    "🔍 Filtres en temps réel (Source, Fenêtre, Auto-refresh)",
    "📱 KPIs cards avec métriques clés",
    "📊 Graphiques interactifs (Plotly)",
    "🔗 Titres d'articles cliquables → Ouverture en nouvel onglet",
    "⚠️ Gestion des quarantaines (Affichage + Approbation)",
    "🔄 Rafraîchissement auto configurable"
])

# Slide 11: Migration & Deployment
add_content_slide(prs, "Déploiement & Migrations", [
    "Docker Compose orchestration:",
    "  • postgres (port 5432) → Data Warehouse",
    "  • streamlit (port 8501) → Dashboard",
    "  • ingest (background) → Scraper cyclique",
    "",
    "Migrations versionnées:",
    "  • 20260507_add_canonical_and_quarantine_audit.sql",
    "  • 20260508_backfill_quarantine_audit.sql (Idempotent)"
])

# Slide 12: Metrics & Results
add_content_slide(prs, "Résultats et Métriques", [
    "📊 Total articles ingérés: 59 uniques",
    "🔄 Distribution par source:",
    "   • BBC News: 22 articles (37%)",
    "   • France24: 20 articles (34%) [NOUVEAU]",
    "   • Hespress: 15 articles (25%)",
    "   • Reuters: 2 articles (3%)",
    "⏱️ Temps ingestion: <1 seconde par cycle"
])

# Slide 13: Next Steps
add_content_slide(prs, "Évolutions Possibles", [
    "🔧 Fine-tuning des heuristiques avec données labelisées",
    "🤖 NLP pour extraction d'entités et thèmes",
    "📤 Export CKAN/Data Portal compatible",
    "📊 Dashboards sectoriels (Tech, Politique, etc.)",
    "⏰ Purge programmée des quartantaines anciennes",
    "🔐 Authentification et RBAC pour production",
    "📧 Alertes et notifications temps réel"
])
# Slide 14: Médaillon Architecture
add_content_slide(prs, "Architecture Médaillon", [
    "🥉 **Bronze Layer** → Articles bruts en MinIO (JSONL)",
    "  Stockage complet de l'historique sans modification",
    "",
    "🥈 **Silver Layer** → Nettoyage & normalisation",
    "  Suppression HTML, normalisation texte, détection langue",
    "  Score de qualité des données",
    "",
    "🥇 **Gold Layer** → Tables analytiques",
    "  articles_by_source, top_keywords, daily_trends"
])

# Slide 15: MinIO Data Lake
add_content_slide(prs, "Data Lake - MinIO", [
    "☁️ Stockage d'objets distributed (S3-compatible)",
    "📁 Buckets organizés par couche:",
    "  • bronze-articles/ → Articles bruts JSONL",
    "  • silver-articles/ → Articles nettoyés Parquet",
    "  • gold-analytics/ → Agrégations CSV",
    "",
    "🔐 Console UI: http://localhost:9001",
    "⚙️ S3-compatible API pour scalabilité production"
])

# Slide 16: Apache Airflow Orchestration
add_content_slide(prs, "Orchestration - Apache Airflow", [
    "🔄 DAG Principal: mediapulse360_pipeline",
    "⏰ Schedule: Horaire (0 * * * *)",
    "",
    "📋 Tâches orchestrées:",
    "  1. extract_bronze → Scrape RSS, store MinIO",
    "  2. transform_silver → Clean, normalize, Quality Score",
    "  3. load_gold → Aggregate, KPI generation",
    "  4. generate_metrics & data_quality_check",
    "",
    "🖥️ UI: http://localhost:8081 (admin/admin)"
])

# Slide 14: Next Steps
# Slide 17: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BLUE

# Main text
main_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
main_frame = main_box.text_frame
main_frame.word_wrap = True
p = main_frame.paragraphs[0]
p.text = "MediaPulse 360"
p.font.size = Pt(50)
p.font.bold = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

p2 = main_frame.add_paragraph()
p2.text = "\nArchitecture Médaillon complète avec Airflow & MinIO\nPrête pour scalabilité et production"
p2.font.size = Pt(24)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(24)

# Footer
footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
footer_frame = footer_box.text_frame
p3 = footer_frame.paragraphs[0]
p3.text = "✅ Architecture Big Data Complète - Prêt pour GitHub & Rapport Final"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(144, 238, 144)
p3.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save("MediaPulse360_Presentation.pptx")
print("✓ Présentation créée: MediaPulse360_Presentation.pptx")
